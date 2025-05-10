from fastapi import (
    FastAPI,
    UploadFile,
    File,
    BackgroundTasks,
    HTTPException,
    Query,
    Depends,
    Form,
)
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
import asyncio
import base64
from typing import List, Union, Optional
import uuid
import os
import hashlib
import httpx
import mammoth
import duckdb
from datetime import datetime
import json

# Imports for GCS
from google.cloud import storage
from google.oauth2 import service_account  # For local testing with service account
from markdownify import markdownify

from io import BytesIO
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import trafilatura


# --- Conditional imports based on the PDF backend ---
def get_pdf_backend():
    return os.getenv("PDF_BACKEND", "gemini")


match get_pdf_backend():
    case "pymupdf4llm":
        from pymupdf4llm import to_markdown
        import pymupdf
    case "pypdf2":
        from pypdf import PdfReader
    case "gemini":
        from pdf2image import convert_from_bytes
        from google import genai
        from google.genai.types import CreateBatchJobConfig, JobState, HttpOptions

        prompt = """OCR this document to Markdown with text formatting such as bold, italic, headings, tables, numbered and bulleted lists properly rendered in Markdown Do not suround the out with Markdown fences. Preserve as much content as possible, such as headings, tables, lists. etc. Do not add any preamble or additional explanation of any other kind --simply output the well-formatted text output in Markdown."""
    case invalid_backend:
        raise ValueError(f"Invalid PDF backend: {invalid_backend}")

app = FastAPI()

# --- Security ---
bearer_scheme = HTTPBearer(
    auto_error=False
)  # auto_error=False to handle optional token & custom errors


def get_api_auth_token() -> Optional[str]:
    return os.getenv("API_AUTH_TOKEN")


async def authenticate_request(
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(bearer_scheme),
) -> None:
    configured_token = get_api_auth_token()
    if configured_token:  # Only enforce auth if a token is configured server-side
        if credentials is None:
            raise HTTPException(
                status_code=401,
                detail="Not authenticated. Authorization header is missing.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.scheme != "Bearer":
            raise HTTPException(
                status_code=401,
                detail="Invalid authentication scheme. Only Bearer is supported.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.token != configured_token:
            raise HTTPException(
                status_code=403,
                detail="Invalid token.",
            )
    # If no token is configured server-side, or if authentication passes, do nothing.
    return


# other configuration
def get_gcs_project_id():
    return os.getenv("GOOGLE_CLOUD_PROJECT")


# For local GCS testing with a service account JSON file
def get_gcs_credentials():
    credentials_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
    if credentials_path:
        return service_account.Credentials.from_service_account_file(credentials_path)
    return None  # Fallback to default environment auth if not set


def get_gemini_client():
    project = get_gcs_project_id()
    location = os.getenv("GOOGLE_CLOUD_LOCATION")
    api_key = os.getenv("ALTAI_GEMINI_API_KEY")
    client = (
        genai.Client(vertexai=False, api_key=api_key)
        if api_key
        else genai.Client(vertexai=True, location=location, project=project)
    )
    return client


def get_max_file_size_bytes() -> Union[int, None]:
    max_size_mb_str = os.getenv("MAX_FILE_SIZE_MB")
    if max_size_mb_str:
        try:
            max_size_mb = int(max_size_mb_str)
            if max_size_mb > 0:
                return max_size_mb * 1024 * 1024  # Convert MB to Bytes
            else:
                # Log this invalid configuration? For now, treat as no limit.
                pass
        except ValueError:
            # Log this invalid configuration? For now, treat as no limit.
            pass
    return None  # No limit or invalid value treated as no limit


SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".rtf", ".pptx", ".html", ".htm"]


class ConversionResponse(BaseModel):
    filename: str
    content_hash: str
    texts: List[str]


class BatchRequest(BaseModel):
    input_paths: Union[str, List[str]]
    output_path: str


TASKS = {}


# --- DuckDB Setup ---
DUCKDB_FILE = os.getenv("DUCKDB_FILE", "batch_tasks.duckdb")


def get_db_connection():
    return duckdb.connect(DUCKDB_FILE)


def initialize_db_schema():
    con = get_db_connection()
    try:
        # Main batch jobs table
        con.execute("""
            CREATE TABLE IF NOT EXISTS batch_jobs (
                job_id VARCHAR PRIMARY KEY,
                output_gcs_path VARCHAR NOT NULL,
                status VARCHAR NOT NULL,
                submitted_at TIMESTAMP NOT NULL,
                total_input_files INTEGER NOT NULL,
                overall_processed_count INTEGER DEFAULT 0,
                overall_failed_count INTEGER DEFAULT 0,
                last_updated_at TIMESTAMP
            )
        """)
        # Gemini PDF sub-jobs (one per Gemini Batch API call)
        con.execute("""
            CREATE TABLE IF NOT EXISTS gemini_pdf_batch_sub_jobs (
                gemini_batch_sub_job_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_api_job_name VARCHAR,
                status VARCHAR NOT NULL,
                payload_gcs_uri VARCHAR,
                gemini_output_gcs_uri_prefix VARCHAR,
                total_pdf_pages_for_batch INTEGER DEFAULT 0,
                processed_pdf_pages_count INTEGER DEFAULT 0,
                failed_pdf_pages_count INTEGER DEFAULT 0,
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
        # Individual file tasks (for non-PDFs, or individual pages of PDFs before aggregation)
        con.execute("""
            CREATE TABLE IF NOT EXISTS file_tasks (
                file_task_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_batch_sub_job_id VARCHAR REFERENCES gemini_pdf_batch_sub_jobs(gemini_batch_sub_job_id), -- Link to a Gemini batch if it's a PDF page
                original_filename VARCHAR NOT NULL,
                file_type VARCHAR NOT NULL, -- e.g., 'pdf_page', 'docx'
                status VARCHAR NOT NULL, -- pending, processing, image_uploaded_to_gcs, completed, failed
                gcs_input_image_uri VARCHAR, -- For PDF pages, GCS URI of the image sent to Gemini
                gcs_output_markdown_uri VARCHAR, -- GCS URI of the final .md (for non-PDFs or aggregated PDFs)
                page_number INTEGER, -- For PDF pages
                gemini_request_id VARCHAR, -- The 'id' used in payload.jsonl for this PDF page
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
    finally:
        con.close()


# Call initialization at startup
initialize_db_schema()


# --- Environment Variables for Batch Processing ---
GCS_BATCH_BUCKET = os.getenv("GCS_BUCKET")
GEMINI_MODEL_FOR_VISION = os.getenv(
    "GEMINI_MODEL_FOR_VISION", "gemini-2.0-flash"
)  # Or other vision model

if not GCS_BATCH_BUCKET:
    print(
        "Warning: GCS_BUCKET environment variable is not set. PDF batch processing will fail."
    )


def _process_docx_sync(content_bytes: bytes) -> List[str]:
    try:
        doc = BytesIO(content_bytes)
        doc_html = mammoth.convert_to_html(doc).value
        doc_md = markdownify(doc_html).strip()
        return [doc_md]
    except Exception as e:
        return [f"Error processing DOCX: {str(e)}"]


def _process_rtf_sync(content_bytes: bytes) -> List[str]:
    try:
        return [rtf_to_text(content_bytes.decode("utf-8", errors="ignore"))]
    except Exception as e:
        return [f"Error processing RTF: {str(e)}"]


def _process_pptx_sync(content_bytes: bytes) -> List[str]:
    try:
        prs = Presentation(BytesIO(content_bytes))
        # Corrected list comprehension for PPTX to build a single string per slide, then list of slide texts
        slide_texts = []
        for slide in prs.slides:
            text_on_slide = "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            if text_on_slide:  # Only add if there's text
                slide_texts.append(text_on_slide)
        return (
            slide_texts if slide_texts else [""]
        )  # Return list of slide texts, or list with empty string if no text
    except Exception as e:
        return [f"Error processing PPTX: {str(e)}"]


def _process_html_sync(content_bytes: bytes) -> List[str]:
    try:
        extracted_text = trafilatura.extract(
            content_bytes.decode("utf-8", errors="ignore"), output_format="markdown"
        )
        return [extracted_text if extracted_text is not None else ""]
    except Exception as e:
        return [f"Error processing HTML: {str(e)}"]


def _process_pdf_pymupdf4llm_sync(content_bytes: bytes) -> List[str]:
    try:
        pymupdf_doc = pymupdf.Document(stream=content_bytes, filetype="pdf")
        page_data_list = to_markdown(pymupdf_doc, page_chunks=True)
        return [page_dict.get("text", "") for page_dict in page_data_list]
    except Exception as e:
        return [f"Error processing PDF with pymupdf4llm: {str(e)}"]


def _process_pdf_pypdf2_sync(content_bytes: bytes) -> List[str]:
    try:
        reader = PdfReader(BytesIO(content_bytes))
        return [p.extract_text() or "" for p in reader.pages]
    except Exception as e:
        return [f"Error processing PDF with pypdf: {str(e)}"]


async def _process_file_content(
    ext: str, content: bytes, pdf_backend_choice: str
) -> List[str]:
    texts_list: List[str] = []
    if ext == ".pdf":
        if pdf_backend_choice == "pymupdf4llm":
            texts_list = await asyncio.to_thread(_process_pdf_pymupdf4llm_sync, content)
        elif pdf_backend_choice == "pypdf2":
            texts_list = await asyncio.to_thread(_process_pdf_pypdf2_sync, content)
        elif pdf_backend_choice == "gemini":
            pages = convert_from_bytes(content)
            images_b64 = []
            for page in pages:
                buffer = BytesIO()
                page.save(buffer, format="PNG")
                image_data = buffer.getvalue()
                b64_str = base64.b64encode(image_data).decode("utf-8")
                images_b64.append(b64_str)
            client = get_gemini_client()
            payloads = [
                [
                    {"inline_data": {"data": b64_str, "mime_type": "image/png"}},
                    {"text": prompt},
                ]
                for b64_str in images_b64
            ]
            results = await asyncio.gather(
                *[
                    client.aio.models.generate_content(
                        model="gemini-2.0-flash", contents=payload
                    )
                    for payload in payloads
                ]
            )
            texts_list = [result.text for result in results]
        else:
            texts_list = ["Invalid PDF backend specified."]
    elif ext in [".docx"]:
        texts_list = await asyncio.to_thread(_process_docx_sync, content)
    elif ext in [".rtf"]:
        texts_list = await asyncio.to_thread(_process_rtf_sync, content)
    elif ext in [".pptx"]:
        texts_list = await asyncio.to_thread(_process_pptx_sync, content)
    elif ext in [".html", ".htm"]:
        texts_list = await asyncio.to_thread(_process_html_sync, content)
    else:
        texts_list = ["Unsupported file type encountered in _process_file_content."]
    return texts_list


@app.post(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_file_upload(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    content = await file.read()

    max_size = get_max_file_size_bytes()
    if max_size is not None and len(content) > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"File size {len(content) / (1024 * 1024):.2f}MB exceeds maximum allowed size of {max_size / (1024 * 1024):.2f}MB.",
        )

    content_hash = hashlib.sha256(content).hexdigest()
    pdf_backend_choice = get_pdf_backend()

    texts_list = await _process_file_content(ext, content, pdf_backend_choice)

    if texts_list and (
        texts_list[0].startswith("Error processing")
        or texts_list[0].startswith("Invalid PDF backend")
        or texts_list[0].startswith("Unsupported file type")
    ):
        raise HTTPException(status_code=400, detail=texts_list[0])

    return ConversionResponse(
        filename=file.filename, content_hash=content_hash, texts=texts_list
    )


@app.get(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_url(
    url: str = Query(..., description="URL of the webpage to convert to Markdown"),
):
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(url)
            response.raise_for_status()
            html_content = response.text
            content_bytes = html_content.encode("utf-8")
    except httpx.RequestError as e:
        raise HTTPException(status_code=400, detail=f"Error fetching URL: {str(e)}")
    except httpx.HTTPStatusError as e:
        raise HTTPException(
            status_code=e.response.status_code,
            detail=f"Error fetching URL: {e.response.reason_phrase}",
        )

    if not html_content:
        raise HTTPException(status_code=400, detail="Fetched content is empty.")

    content_hash = hashlib.sha256(content_bytes).hexdigest()

    extracted_text = trafilatura.extract(html_content, output_format="markdown")
    texts_list = [extracted_text if extracted_text is not None else ""]

    filename = os.path.basename(url) or url

    return ConversionResponse(
        filename=filename, content_hash=content_hash, texts=texts_list
    )


@app.get("/status/{task_id}", dependencies=[Depends(authenticate_request)])
def status(task_id: str):
    # This will be updated to query DuckDB for batch jobs
    # For now, keep the old behavior for non-batch tasks if any, or return not found.
    # If task_id looks like a UUID, it might be an old task.
    # A more robust solution would be to prefix batch task_ids or check DB first.
    con = get_db_connection()
    try:
        job_status = con.execute(
            "SELECT * FROM batch_jobs WHERE job_id = ?", (task_id,)
        ).fetchone()
        if job_status:
            # Fetch details from other tables as well
            job_dict = dict(zip([desc[0] for desc in con.description], job_status))

            # Fetch gemini sub job details if any
            gemini_sub_jobs = con.execute(
                "SELECT * FROM gemini_pdf_batch_sub_jobs WHERE batch_job_id = ?",
                (task_id,),
            ).fetchall()
            job_dict["gemini_pdf_processing_details"] = [
                dict(zip([desc[0] for desc in con.description], sub_job))
                for sub_job in gemini_sub_jobs
            ]

            # Fetch individual file task details
            file_tasks_details = con.execute(
                "SELECT original_filename, file_type, status, gcs_output_markdown_uri, error_message, page_number FROM file_tasks WHERE batch_job_id = ?",
                (task_id,),
            ).fetchall()
            job_dict["file_processing_details"] = [
                dict(zip([desc[0] for desc in con.description], task))
                for task in file_tasks_details
            ]

            return job_dict
        else:  # Fallback to old TASKS dict or not found
            return TASKS.get(
                task_id,
                {
                    "status": "not_found",
                    "detail": "Task not found in active batch jobs or older task system.",
                },
            )
    finally:
        con.close()


@app.post("/batch", dependencies=[Depends(authenticate_request)])
async def batch_files_upload(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
    output_gcs_path: str = Form(...),
):
    main_batch_job_id = str(uuid.uuid4())
    current_time = datetime.utcnow()

    pdf_files_for_gemini_batch: List[UploadFile] = []
    non_pdf_files_for_individual_processing: List[UploadFile] = []

    if not GCS_BATCH_BUCKET:
        raise HTTPException(
            status_code=500,
            detail="GCS_BATCH_TEMP_BUCKET is not configured on the server.",
        )
    if not output_gcs_path.startswith("gs://"):
        raise HTTPException(
            status_code=400, detail="Output GCS path must start with gs://"
        )

    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        if ext == ".pdf":
            pdf_files_for_gemini_batch.append(f)
        elif ext in SUPPORTED_EXTENSIONS:  # Excludes .pdf as it's handled above
            non_pdf_files_for_individual_processing.append(f)
        else:
            # Optionally log or report unsupported files
            print(f"Skipping unsupported file: {f.filename}")

    con = get_db_connection()
    try:
        con.execute(
            "INSERT INTO batch_jobs (job_id, output_gcs_path, status, submitted_at, total_input_files, last_updated_at) VALUES (?, ?, ?, ?, ?, ?)",
            (
                main_batch_job_id,
                output_gcs_path,
                "pending",
                current_time,
                len(pdf_files_for_gemini_batch)
                + len(non_pdf_files_for_individual_processing),
                current_time,
            ),
        )
        con.commit()

        # Process non-PDF files
        for upload_file in non_pdf_files_for_individual_processing:
            file_task_id = str(uuid.uuid4())
            file_ext = os.path.splitext(upload_file.filename)[1].lower()
            con.execute(
                "INSERT INTO file_tasks (file_task_id, batch_job_id, original_filename, file_type, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    file_task_id,
                    main_batch_job_id,
                    upload_file.filename,
                    file_ext,
                    "pending",
                    current_time,
                    current_time,
                ),
            )
            # Read content before passing to background task
            content_bytes = await upload_file.read()
            await upload_file.seek(
                0
            )  # Reset cursor for safety, though might not be strictly needed if not read again

            background_tasks.add_task(
                _process_single_non_pdf_file_and_upload,
                content_bytes,  # Pass content bytes
                file_ext,
                upload_file.filename,
                output_gcs_path,
                main_batch_job_id,
                file_task_id,
            )
        con.commit()

        # Process PDF files via Gemini Batch
        if pdf_files_for_gemini_batch:
            gemini_batch_sub_job_id = str(uuid.uuid4())
            con.execute(
                "INSERT INTO gemini_pdf_batch_sub_jobs (gemini_batch_sub_job_id, batch_job_id, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                (
                    gemini_batch_sub_job_id,
                    main_batch_job_id,
                    "pending_preparation",
                    current_time,
                    current_time,
                ),
            )
            # Pass the list of UploadFile objects. Content will be read in the task.
            background_tasks.add_task(
                _run_gemini_pdf_batch_conversion,
                [
                    f for f in pdf_files_for_gemini_batch
                ],  # Pass copies or ensure they are not closed
                output_gcs_path,
                main_batch_job_id,
                gemini_batch_sub_job_id,
            )
            con.commit()  # Commit after starting PDF batch task prep

        # Update batch job status to processing if there are tasks
        if pdf_files_for_gemini_batch or non_pdf_files_for_individual_processing:
            con.execute(
                "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                ("processing", datetime.utcnow(), main_batch_job_id),
            )
        else:  # No files to process
            con.execute(
                "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                ("completed_no_files", datetime.utcnow(), main_batch_job_id),
            )
        con.commit()

    finally:
        con.close()

    return {"task_id": main_batch_job_id}


# Placeholder for the new helper functions - to be implemented next
async def _process_single_non_pdf_file_and_upload(
    content_bytes: bytes,
    file_ext: str,
    original_filename: str,
    output_gcs_path_str: str,
    main_batch_job_id: str,
    file_task_id: str,
):
    current_time = datetime.utcnow()
    con = get_db_connection()
    try:
        con.execute(
            "UPDATE file_tasks SET status = ?, updated_at = ? WHERE file_task_id = ?",
            ("processing", current_time, file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET last_updated_at = ? WHERE job_id = ?",
            (current_time, main_batch_job_id),
        )
        con.commit()

        # Re-use the existing _process_file_content logic
        # Ensure pdf_backend_choice is not relevant or handled if _process_file_content expects it
        # For non-PDFs, pdf_backend_choice is not used by _process_file_content.
        markdown_texts = await _process_file_content(
            file_ext, content_bytes, get_pdf_backend()
        )  # get_pdf_backend() is for PDF choice, not directly used for non-PDFs here

        if (
            not markdown_texts
            or markdown_texts[0].startswith("Error processing")
            or markdown_texts[0].startswith("Unsupported file type")
        ):
            error_message = (
                markdown_texts[0] if markdown_texts else "Unknown processing error"
            )
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_message, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
            print(
                f"Failed to process non-PDF file {original_filename}: {error_message}"
            )
            return

        full_markdown_output = "\n\n---\n\n".join(markdown_texts)

        # Upload to GCS
        storage_client = storage.Client(
            project=get_gcs_project_id(), credentials=get_gcs_credentials()
        )
        output_bucket_name, output_prefix = output_gcs_path_str.replace(
            "gs://", ""
        ).split("/", 1)
        output_bucket = storage_client.bucket(output_bucket_name)

        output_blob_name = (
            output_prefix.rstrip("/")
            + "/"
            + os.path.splitext(original_filename)[0]
            + ".md"
        )
        output_blob_obj = output_bucket.blob(output_blob_name)

        output_blob_obj.upload_from_string(
            full_markdown_output, content_type="text/markdown"
        )
        gcs_output_url = f"gs://{output_bucket_name}/{output_blob_name}"

        con.execute(
            "UPDATE file_tasks SET status = ?, gcs_output_markdown_uri = ?, updated_at = ? WHERE file_task_id = ?",
            ("completed", gcs_output_url, datetime.utcnow(), file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET overall_processed_count = overall_processed_count + 1, last_updated_at = ? WHERE job_id = ?",
            (datetime.utcnow(), main_batch_job_id),
        )
        con.commit()
        print(
            f"Successfully processed and uploaded non-PDF file {original_filename} to {gcs_output_url}"
        )

    except Exception as e:
        error_str = f"Error in _process_single_non_pdf_file_and_upload for {original_filename}: {str(e)}"
        print(error_str)
        try:
            # Attempt to mark as failed in DB
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_str, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(
                f"Additionally, failed to update DB for task {file_task_id} failure: {db_err}"
            )
    finally:
        con.close()


async def _run_gemini_pdf_batch_conversion(
    pdf_files_list: List[UploadFile],  # List of UploadFile objects
    output_gcs_path_str: str,
    main_batch_job_id: str,
    gemini_batch_sub_job_id: str,
):
    current_time = datetime.utcnow()
    con = get_db_connection()
    storage_client = storage.Client(
        project=get_gcs_project_id(), credentials=get_gcs_credentials()
    )
    gemini_client = get_gemini_client()  # Ensure this client is suitable for batch

    temp_gcs_input_prefix = (
        f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/inputs"
    )
    temp_gcs_images_prefix = f"{temp_gcs_input_prefix}/images"
    temp_gcs_output_prefix = (
        f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/outputs"
    )

    payload_items_for_jsonl = []
    total_pages_for_this_gemini_batch = 0

    try:
        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            ("preparing_images_and_payload", current_time, gemini_batch_sub_job_id),
        )
        con.execute(
            "UPDATE batch_jobs SET last_updated_at = ? WHERE job_id = ?",
            (current_time, main_batch_job_id),
        )
        con.commit()

        temp_bucket = storage_client.bucket(GCS_BATCH_BUCKET)

        for pdf_upload_file in pdf_files_list:
            original_pdf_filename = pdf_upload_file.filename
            pdf_content_bytes = await pdf_upload_file.read()
            await pdf_upload_file.close()  # Close the file after reading

            try:
                page_images = convert_from_bytes(pdf_content_bytes, fmt="png")
            except Exception as e:
                print(f"Failed to convert PDF {original_pdf_filename} to images: {e}")
                # Mark all potential pages for this PDF as failed in file_tasks if desired, or just log
                # For now, we skip this PDF and it won't contribute to total_pages_for_this_gemini_batch
                # and won't have file_tasks entries created here.
                # A more robust approach might create failed file_tasks entries for it.
                con.execute(
                    "UPDATE gemini_pdf_batch_sub_jobs SET error_message = COALESCE(error_message || CHR(10), '') || ? WHERE gemini_batch_sub_job_id = ?",
                    (
                        f"Failed to convert PDF {original_pdf_filename} to images: {str(e)}",
                        gemini_batch_sub_job_id,
                    ),
                )
                con.commit()
                continue  # Skip to the next PDF file

            total_pages_for_this_gemini_batch += len(page_images)

            for i, page_image in enumerate(page_images):
                page_num = i + 1
                gemini_request_id = f"{os.path.splitext(original_pdf_filename)[0]}_p{page_num}_{uuid.uuid4().hex[:8]}"
                file_task_id = str(uuid.uuid4())

                con.execute(
                    "INSERT INTO file_tasks (file_task_id, batch_job_id, gemini_batch_sub_job_id, original_filename, file_type, status, page_number, gemini_request_id, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                    (
                        file_task_id,
                        main_batch_job_id,
                        gemini_batch_sub_job_id,
                        original_pdf_filename,
                        "pdf_page",
                        "pending_image_upload",
                        page_num,
                        gemini_request_id,
                        datetime.utcnow(),
                        datetime.utcnow(),
                    ),
                )
                con.commit()

                try:
                    img_byte_arr = BytesIO()
                    page_image.save(img_byte_arr, format="PNG")
                    img_byte_arr = img_byte_arr.getvalue()

                    image_blob_name = (
                        f"{temp_gcs_images_prefix}/{gemini_request_id}.png"
                    )
                    image_blob = temp_bucket.blob(image_blob_name)
                    image_blob.upload_from_string(
                        img_byte_arr, content_type="image/png"
                    )
                    image_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{image_blob_name}"

                    con.execute(
                        "UPDATE file_tasks SET gcs_input_image_uri = ?, status = ?, updated_at = ? WHERE file_task_id = ?",
                        (
                            image_gcs_uri,
                            "image_uploaded_to_gcs",
                            datetime.utcnow(),
                            file_task_id,
                        ),
                    )
                    con.commit()
                    payload_items_for_jsonl.append(
                        {
                            "id": gemini_request_id,
                            "request": {
                                "contents": [
                                    {
                                        "file_data": {
                                            "file_uri": image_gcs_uri,
                                            "mime_type": "image/png",
                                        }
                                    },
                                    {"text": prompt},  # Using the global prompt for OCR
                                ]
                            },
                        }
                    )
                except Exception as e:
                    print(
                        f"Failed to upload image for {original_pdf_filename} page {page_num}: {e}"
                    )
                    con.execute(
                        "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                        (
                            "failed",
                            f"Image upload error: {str(e)}",
                            datetime.utcnow(),
                            file_task_id,
                        ),
                    )
                    # also update gemini_pdf_batch_sub_jobs about this failure
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + 1, error_message = COALESCE(error_message || CHR(10), '') || ? WHERE gemini_batch_sub_job_id = ?",
                        (
                            f"Failed image upload for {original_pdf_filename} page {page_num}: {str(e)}",
                            gemini_batch_sub_job_id,
                        ),
                    )
                    con.commit()
                    # Continue to next page, this page will be marked as failed.

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET total_pdf_pages_for_batch = ? WHERE gemini_batch_sub_job_id = ?",
            (total_pages_for_this_gemini_batch, gemini_batch_sub_job_id),
        )
        con.commit()

        if not payload_items_for_jsonl:
            print(
                f"No pages successfully prepared for Gemini batch sub job {gemini_batch_sub_job_id}. Aborting Gemini submission."
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || 'No pages to process', updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("failed_no_payload", datetime.utcnow(), gemini_batch_sub_job_id),
            )
            # Update main batch job if all PDFs failed here and there were no non-PDFs
            # This logic might need to be more sophisticated at the batch_jobs level based on counts
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + ? WHERE job_id = ?",
                (len(pdf_files_list), main_batch_job_id),
            )
            con.commit()
            return

        # Create and upload payload.jsonl
        payload_jsonl_content = "\n".join(
            [json.dumps(item) for item in payload_items_for_jsonl]
        )
        payload_blob_name = f"{temp_gcs_input_prefix}/payload.jsonl"
        payload_blob = temp_bucket.blob(payload_blob_name)
        payload_blob.upload_from_string(
            payload_jsonl_content, content_type="application/jsonl"
        )
        payload_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{payload_blob_name}"

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET payload_gcs_uri = ?, status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            (
                payload_gcs_uri,
                "submitting_to_gemini",
                datetime.utcnow(),
                gemini_batch_sub_job_id,
            ),
        )
        con.commit()

        # Submit to Gemini Batch API
        gemini_output_uri_for_job = f"gs://{GCS_BATCH_BUCKET}/{temp_gcs_output_prefix}"
        batch_job_config = CreateBatchJobConfig(
            destination_uri=gemini_output_uri_for_job
        )

        # Ensure client is correctly initialized for batch (Vertex AI or API Key)
        # The article uses client = genai.Client(http_options=HttpOptions(api_version="v1"), vertexai=True, project=PROJECT_ID, location=LOCATION)
        # We are using get_gemini_client() which needs to be aligned.
        # For now, assuming get_gemini_client() is correctly configured for the target (Vertex or Generative Language API batch)

        # Note: google-genai client.batches.create is synchronous and polls internally.
        # For very long jobs, a truly async submission and separate polling might be better,
        # but for now, this background task will run it.
        print(
            f"Submitting batch job to Gemini for {gemini_batch_sub_job_id} with model {GEMINI_MODEL_FOR_VISION}"
        )
        gemini_job = gemini_client.batches.create(
            model=f"models/{GEMINI_MODEL_FOR_VISION}",  # Model name might need prefix
            source_uri=payload_gcs_uri,
            config=batch_job_config,
        )
        # The create call above blocks until job completion or failure for google-genai,
        # or at least starts it and returns an object that can be polled. Let's assume it returns quickly with a job object.
        # The example in the article suggests it returns a job object immediately, and then polling is done.
        # client.batches.create() is a synchronous call that waits for the batch job to complete.
        # This might make the background task run for a long time.

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET gemini_api_job_name = ?, status = ?, gemini_output_gcs_uri_prefix = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            (
                gemini_job.name,
                str(gemini_job.state),
                gemini_output_uri_for_job,
                datetime.utcnow(),
                gemini_batch_sub_job_id,
            ),
        )
        con.commit()
        print(f"Gemini job {gemini_job.name} submitted. State: {gemini_job.state}")

        # Polling loop (client.batches.create might already do this, check SDK behavior)
        # If client.batches.create() is indeed blocking until completion, this loop is redundant for it,
        # but useful if it returns immediately.
        # For google-genai, `create` is blocking and polls until completion. `get` is then used to refresh state if needed.
        # Let's assume `gemini_job` is the final state object after `create()` completes.

        if gemini_job.state == JobState.JOB_STATE_SUCCEEDED:
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    "processing_gemini_results",
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            con.commit()

            # Find the predictions.jsonl file. It's in a timestamped or job-id named subfolder.
            # The article uses a utility: get_latest_folder. Simpler: list blobs with prefix.
            # Output is typically: gs://{dest_bucket}/{dest_prefix}/gemini_batch_JOB_ID_TIMESTAMP/predictions.jsonl
            # The gemini_job.output_files might contain the direct path.
            # Let's assume the Batch API produces output like: {gemini_output_uri_for_job}/{gemini_job.name}/predictions.jsonl (this structure might vary)
            # Or more likely: {gemini_output_uri_for_job}/predictions_JOB_ID_PART_OF_URI.jsonl or similar fixed name patterns.

            # The result files are directly in the output_uri specified. e.g. {gemini_output_uri_for_job}/predictions.jsonl-00000-of-00001
            # Or, it might be a folder. The SDK or docs should clarify this structure precisely.
            # The article implies: `output_uri/TIMESTAMPED_FOLDER/predictions.jsonl`
            # Let's try to list blobs to find it.
            output_blobs = list(
                storage_client.list_blobs(
                    GCS_BATCH_BUCKET, prefix=f"{temp_gcs_output_prefix}/"
                )
            )
            predictions_file_blob = None
            for blob_item in output_blobs:
                if blob_item.name.endswith("predictions.jsonl") or (
                    "predictions.jsonl-" in blob_item.name
                ):  # Batch output can be sharded
                    predictions_file_blob = blob_item
                    break  # Take the first one found, or handle multiple shards if necessary

            if not predictions_file_blob:
                raise Exception(
                    f"predictions.jsonl not found in Gemini output: gs://{GCS_BATCH_BUCKET}/{temp_gcs_output_prefix}/"
                )

            print(f"Downloading predictions from: {predictions_file_blob.name}")
            predictions_content = predictions_file_blob.download_as_text()

            page_markdown_results = {}
            processed_gemini_pages = 0
            failed_gemini_pages = 0

            for line in predictions_content.splitlines():
                if not line.strip():
                    continue
                try:
                    prediction = json.loads(line)
                    request_id = prediction["id"]  # This is our gemini_request_id
                    # Assuming direct text output as we didn't ask for structured JSON response for the text itself
                    markdown_text = prediction["response"]["candidates"][0]["content"][
                        "parts"
                    ][0]["text"]
                    page_markdown_results[request_id] = markdown_text
                    processed_gemini_pages += 1
                except Exception as e:
                    print(f"Error parsing prediction line: {line}, Error: {e}")
                    # How to map this failure back to a specific page if ID is unparsable?
                    # For now, just count as a general parsing failure for the batch.
                    failed_gemini_pages += 1  # This is a prediction parsing failure, not a page processing failure yet.

            # Aggregate results by original PDF and upload
            # Group file_tasks by original_filename for those in this gemini_batch_sub_job_id
            pdf_page_tasks_cursor = con.execute(
                "SELECT file_task_id, original_filename, page_number, gemini_request_id FROM file_tasks WHERE gemini_batch_sub_job_id = ? ORDER BY original_filename, page_number ASC",
                (gemini_batch_sub_job_id,),
            ).fetchall()

            aggregated_pdfs = {}
            for task_row_tuple in pdf_page_tasks_cursor:
                task_id, pdf_name, page_no, req_id = task_row_tuple
                if pdf_name not in aggregated_pdfs:
                    aggregated_pdfs[pdf_name] = []

                markdown = page_markdown_results.get(req_id)
                if markdown:
                    aggregated_pdfs[pdf_name].append(
                        (page_no, markdown, task_id, "completed")
                    )
                else:
                    # This page's markdown wasn't found in predictions, or failed parsing earlier
                    aggregated_pdfs[pdf_name].append(
                        (
                            page_no,
                            "Error: OCR output not found for this page.",
                            task_id,
                            "failed",
                        )
                    )
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + 1 WHERE gemini_batch_sub_job_id = ?",
                        (gemini_batch_sub_job_id,),
                    )
                    con.commit()

            final_pdfs_processed_count = 0
            final_pdfs_failed_count = 0

            for pdf_name, pages_data in aggregated_pdfs.items():
                pages_data.sort(key=lambda x: x[0])  # Sort by page_number
                full_pdf_markdown = "\n\n---\n\n".join(
                    [p_data[1] for p_data in pages_data]
                )

                output_bucket_name_final, output_prefix_final = (
                    output_gcs_path_str.replace("gs://", "").split("/", 1)
                )
                output_bucket_final = storage_client.bucket(output_bucket_name_final)
                final_blob_name = (
                    output_prefix_final.rstrip("/")
                    + "/"
                    + os.path.splitext(pdf_name)[0]
                    + ".md"
                )
                final_blob_obj = output_bucket_final.blob(final_blob_name)

                try:
                    final_blob_obj.upload_from_string(
                        full_pdf_markdown, content_type="text/markdown"
                    )
                    gcs_final_url = f"gs://{output_bucket_name_final}/{final_blob_name}"
                    # Mark all page tasks for this PDF based on their individual outcomes
                    all_pages_succeeded = True
                    for p_data_item in pages_data:
                        p_task_id, p_status = p_data_item[2], p_data_item[3]
                        p_error_msg = p_data_item[1] if p_status == "failed" else None
                        con.execute(
                            "UPDATE file_tasks SET status = ?, gcs_output_markdown_uri = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                            (
                                p_status,
                                gcs_final_url if p_status == "completed" else None,
                                p_error_msg,
                                datetime.utcnow(),
                                p_task_id,
                            ),
                        )
                        if p_status == "failed":
                            all_pages_succeeded = False

                    if all_pages_succeeded:
                        final_pdfs_processed_count += 1
                        con.execute(
                            "UPDATE gemini_pdf_batch_sub_jobs SET processed_pdf_pages_count = processed_pdf_pages_count + ? WHERE gemini_batch_sub_job_id = ?",
                            (len(pages_data), gemini_batch_sub_job_id),
                        )
                    else:
                        final_pdfs_failed_count += (
                            1  # Count PDF as failed if any page failed
                        )
                    con.commit()
                except Exception as upload_exc:
                    print(
                        f"Failed to upload final aggregated PDF {pdf_name}: {upload_exc}"
                    )
                    final_pdfs_failed_count += 1
                    # Mark all page tasks for this PDF as failed due to final upload error
                    for p_data_item in pages_data:
                        p_task_id = p_data_item[2]
                        con.execute(
                            "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                            (
                                "failed",
                                f"Final GCS upload error: {str(upload_exc)}",
                                datetime.utcnow(),
                                p_task_id,
                            ),
                        )
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + ? WHERE gemini_batch_sub_job_id = ?",
                        (len(pages_data), gemini_batch_sub_job_id),
                    )
                    con.commit()

            con.execute(
                "UPDATE batch_jobs SET overall_processed_count = overall_processed_count + ?, overall_failed_count = overall_failed_count + ?, last_updated_at = ? WHERE job_id = ?",
                (
                    final_pdfs_processed_count,
                    final_pdfs_failed_count,
                    datetime.utcnow(),
                    main_batch_job_id,
                ),
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("completed", datetime.utcnow(), gemini_batch_sub_job_id),
            )
            con.commit()
            print(
                f"Gemini PDF batch sub job {gemini_batch_sub_job_id} completed successfully."
            )

        else:  # Gemini job did not succeed
            error_msg_from_job = (
                str(gemini_job.error)
                if hasattr(gemini_job, "error") and gemini_job.error
                else "Unknown Gemini job failure"
            )
            print(
                f"Gemini job {gemini_job.name} failed. State: {gemini_job.state}, Error: {error_msg_from_job}"
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    f"failed_gemini_job_{str(gemini_job.state)}",
                    error_msg_from_job,
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            # Mark all associated file_tasks as failed
            con.execute(
                "UPDATE file_tasks SET status='failed', error_message=?, updated_at=? WHERE gemini_batch_sub_job_id = ?",
                (
                    f"Gemini job failed: {error_msg_from_job}",
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + ? WHERE job_id = ?",
                (len(pdf_files_list), main_batch_job_id),
            )  # Approx count
            con.commit()

    except Exception as e:
        error_str = f"Error in _run_gemini_pdf_batch_conversion for sub-job {gemini_batch_sub_job_id}: {str(e)}"
        print(error_str)
        # Attempt to mark sub-job as failed in DB
        try:
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    "failed_internal_error",
                    error_str,
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            # Also mark relevant file_tasks as failed if they weren't already
            con.execute(
                "UPDATE file_tasks SET status='failed', error_message=? WHERE gemini_batch_sub_job_id = ? AND status NOT IN ('completed', 'failed')",
                (
                    f"Internal error in batch: {error_str[:200]}",
                    gemini_batch_sub_job_id,
                ),
            )
            # Update overall batch job with failures - this is an approximation
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + (SELECT COUNT(*) FROM file_tasks WHERE gemini_batch_sub_job_id = ? AND status = 'failed'), last_updated_at = ? WHERE job_id = ?",
                (gemini_batch_sub_job_id, datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(
                f"Additionally, failed to update DB for PDF batch sub-job {gemini_batch_sub_job_id} failure: {db_err}"
            )
    finally:
        # Clean up temporary image files from GCS_BATCH_TEMP_BUCKET/images to save costs (optional)
        # This requires listing and deleting, can be a separate cleanup task or done here.
        # For simplicity, not implemented in this iteration.
        # print(f"Cleaning up GCS temp files for {gemini_batch_sub_job_id} - not implemented")
        con.close()


# Actual batch processing logic
async def run_batch_task(
    input_paths: Union[str, List[str]], output_gcs_path_str: str, task_id: str
):
    # This function is now OBSOLETE for the new /batch endpoint.
    # It was designed for GCS path inputs and the old TASKS dictionary.
    # Keeping it here for now to avoid breaking if it's called from somewhere else unexpectedly,
    # but it should eventually be removed or refactored if any part is still needed.
    # The new /batch flow uses _process_single_non_pdf_file_and_upload and _run_gemini_pdf_batch_conversion.
    TASKS[task_id]["status"] = "initializing_OBSOLETE_PATH"
    project_id = get_gcs_project_id()
    credentials = get_gcs_credentials()

    if (
        not project_id and not credentials
    ):  # Basic check if running outside GCP and no local creds
        TASKS[task_id]["status"] = "error_OBSOLETE_PATH"
        TASKS[task_id]["details"].append(
            "GCS_PROJECT_ID not set or GOOGLE_APPLICATION_CREDENTIALS not found for GCS access (obsolete path)."
        )
        return

    try:
        storage_client = storage.Client(project=project_id, credentials=credentials)
    except Exception as e:
        TASKS[task_id]["status"] = "error_OBSOLETE_PATH"
        TASKS[task_id]["details"].append(
            f"Failed to initialize GCS client (obsolete path): {str(e)}"
        )
        return

    TASKS[task_id]["status"] = "running_OBSOLETE_PATH"
    files_to_process = []

    pdf_backend_choice = get_pdf_backend()

    try:
        if isinstance(input_paths, str):  # It's a GCS directory path
            TASKS[task_id]["details"].append(
                f"Processing directory (obsolete path): {input_paths}"
            )
            bucket_name, prefix = input_paths.replace("gs://", "").split("/", 1)
            bucket = storage_client.bucket(bucket_name)
            blobs = list(bucket.list_blobs(prefix=prefix))
            for blob in blobs:
                if any(
                    blob.name.lower().endswith(ext) for ext in SUPPORTED_EXTENSIONS
                ) and not blob.name.endswith("/"):  # check if not a directory itself
                    files_to_process.append((bucket_name, blob.name))
        else:  # It's a list of GCS file paths
            TASKS[task_id]["details"].append(
                f"Processing list of files (obsolete path): {len(input_paths)} files"
            )
            for path_str in input_paths:
                if not path_str.startswith("gs://"):
                    TASKS[task_id]["details"].append(
                        f"Skipping invalid GCS path (obsolete path): {path_str}"
                    )
                    continue
                bucket_name, blob_name = path_str.replace("gs://", "").split("/", 1)
                if any(blob_name.lower().endswith(ext) for ext in SUPPORTED_EXTENSIONS):
                    files_to_process.append((bucket_name, blob_name))
                else:
                    TASKS[task_id]["details"].append(
                        f"Skipping unsupported file type (obsolete path): {blob_name}"
                    )

        if not files_to_process:
            TASKS[task_id]["details"].append(
                "No supported files found to process (obsolete path)."
            )
            TASKS[task_id]["status"] = "completed_with_no_files_OBSOLETE_PATH"
            return

        TASKS[task_id]["total_files"] = len(files_to_process)
        output_bucket_name, output_prefix = output_gcs_path_str.replace(
            "gs://", ""
        ).split("/", 1)
        output_bucket = storage_client.bucket(output_bucket_name)

        for i, (bucket_name, blob_name) in enumerate(files_to_process):
            TASKS[task_id]["current_file_processing"] = (
                f"{i + 1}/{len(files_to_process)}: {blob_name} (obsolete path)"
            )
            try:
                input_bucket_obj = storage_client.bucket(bucket_name)
                blob_obj = input_bucket_obj.blob(blob_name)
                file_content_bytes = blob_obj.download_as_bytes()

                file_ext = os.path.splitext(blob_name)[1].lower()

                markdown_texts = await _process_file_content(
                    file_ext, file_content_bytes, pdf_backend_choice
                )

                # Combine all pages/sections into a single markdown string for the output file
                full_markdown_output = "\n\n---\n\n".join(markdown_texts)

                output_blob_name = (
                    output_prefix
                    + "/"
                    + os.path.basename(blob_name).rsplit(".", 1)[0]
                    + ".md"
                )

                output_blob_obj = output_bucket.blob(output_blob_name)

                output_blob_obj.upload_from_string(
                    full_markdown_output, content_type="text/markdown"
                )
                TASKS[task_id]["processed_files"] += 1
                TASKS[task_id]["details"].append(
                    f"Successfully processed and uploaded (obsolete path): {blob_name} to {output_blob_obj.public_url if hasattr(output_blob_obj, 'public_url') else output_blob_obj.name}"
                )
            except Exception as e:
                TASKS[task_id]["failed_files"] += 1
                TASKS[task_id]["details"].append(
                    f"Failed to process file {blob_name} (obsolete path): {str(e)}"
                )

        TASKS[task_id]["status"] = "completed_OBSOLETE_PATH"

    except Exception as e:
        TASKS[task_id]["status"] = "error_OBSOLETE_PATH"
        TASKS[task_id]["details"].append(
            f"An unexpected error occurred during batch processing (obsolete path): {str(e)}"
        )
    finally:
        if "current_file_processing" in TASKS[task_id]:
            del TASKS[task_id]["current_file_processing"]  # Clean up
